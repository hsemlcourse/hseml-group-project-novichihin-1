"""Generate CP3 defense presentation (PPTX) programmatically."""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT = RGBColor(0x00, 0x96, 0xD6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)
GREEN = RGBColor(0x00, 0xC8, 0x53)
RED = RGBColor(0xE0, 0x40, 0x40)
ORANGE = RGBColor(0xFF, 0x9F, 0x43)


def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, font_name="Calibri", spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        p.level = 0
    return tf


def add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = "Calibri"
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
                paragraph.font.name = "Calibri"
                paragraph.alignment = PP_ALIGN.CENTER
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    return table


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ================================================================
    # SLIDE 1 — Title
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)

    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                 "Предсказание прибыли\ndigital-рекламной кампании",
                 font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.6),
                 "ML-регрессия: таргетинг + креатив + аукцион → profit",
                 font_size=22, color=ACCENT, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(1.2),
                 "Новичихин Степан Алексеевич  |  БИВ232\nМИЭМ НИУ ВШЭ  •  2026",
                 font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # ================================================================
    # SLIDE 2 — Задача и данные
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
                 "Задача и данные", font_size=32, bold=True, color=ACCENT)

    items = [
        "Цель: предсказать profit кампании до/во время запуска → оптимизация бюджета",
        "Задача: регрессия (profit = revenue − ad_spend, может быть < 0)",
        "Датасет: Kaggle — 10 000 строк, 41 колонка (после очистки ~35 фичей)",
        "Метрика: RMSE (primary) — штрафует крупные ошибки сильнее. MAE, R² — вспомогательные",
        "MAPE / RMSLE не применимы — profit принимает отрицательные значения",
        "Контроль leakage: revenue, ad_spend, ROAS, CPA, CPC удалены (реконструируют таргет)",
        "Сплит: time-based 70/15/15 — модель не «видит будущее» при обучении",
    ]
    add_bullet_list(slide, Inches(0.7), Inches(1.3), Inches(11.5), Inches(5.5),
                    items, font_size=18)

    # ================================================================
    # SLIDE 3 — Признаки и FE
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
                 "Признаки и Feature Engineering", font_size=32, bold=True, color=ACCENT)

    headers = ["Группа", "Примеры", "Обработка"]
    rows = [
        ["Категориальные (17)", "platform, device_type, industry_vertical, income_bracket", "OneHotEncoder"],
        ["Числовые (13)", "impressions, clicks, conversions, quality_score, CTR", "Imputer + StandardScaler"],
        ["Булевы (2)", "has_call_to_action, retargeting_flag", "Imputer (most_frequent)"],
        ["Дата (1 → 5)", "start_date → month, dayofweek, weekofyear, is_weekend, dayofyear", "Imputer + Scaler"],
        ["Удалены (leakage)", "revenue, ad_spend, ROAS, CPA, CPC, actual_cpc", "Dropped (= profit)"],
    ]
    add_table(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(3.2), headers, rows,
              col_widths=[Inches(2.5), Inches(6.0), Inches(3.8)])

    add_text_box(slide, Inches(0.7), Inches(5.0), Inches(11.5), Inches(1.5),
                 "Всё внутри sklearn Pipeline — статистики учатся только на train.\n"
                 "Baseline (02_baseline.ipynb): без календарных фичей — проверяем вклад date FE.",
                 font_size=16, color=LIGHT_GRAY)

    # ================================================================
    # SLIDE 4 — Baseline (CP1)
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
                 "Baseline и эксперименты CP1", font_size=32, bold=True, color=ACCENT)

    headers = ["Модель", "RMSE (val)", "RMSE (test)", "MAE (test)", "R² (test)"]
    rows = [
        ["Dummy (median)", "73 946", "127 486", "32 116", "−0.057"],
        ["LinearRegression*", "31 564", "72 691", "21 931", "0.656"],
        ["Ridge (α=100)", "30 739", "72 869", "21 378", "0.655"],
        ["Lasso (α=10)", "31 517", "72 702", "21 882", "0.656"],
        ["KNN (k=5)", "40 106", "79 389", "21 424", "0.590"],
        ["RandomForest (300)", "41 051", "69 462", "14 869", "0.686"],
    ]
    add_table(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(3.8), headers, rows,
              col_widths=[Inches(3.0), Inches(2.0), Inches(2.2), Inches(2.2), Inches(2.9)])

    # highlight the winner
    add_text_box(slide, Inches(0.7), Inches(5.5), Inches(11.5), Inches(1.0),
                 "* LinearRegression — baseline без calendar FE.  "
                 "RandomForest — лидер CP1 по RMSE test и MAE test.",
                 font_size=16, color=LIGHT_GRAY)

    # ================================================================
    # SLIDE 5 — Experiments CP2
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
                 "Эксперименты CP2: бустинг, стэкинг, PCA", font_size=32, bold=True, color=ACCENT)

    headers = ["Эксперимент", "Модель", "RMSE (val)", "RMSE (test)", "R² (test)"]
    rows = [
        ["exp01_rf", "RandomForest", "40 084", "69 441", "0.686"],
        ["exp02_boost_tuned", "HistGradientBoosting", "30 876", "79 368", "0.590"],
        ["exp03_tree_tuned ★", "GradientBoosting", "26 778", "64 019", "0.733"],
        ["exp04_stack", "Stacking (RF+Boost→Ridge)", "34 734", "72 138", "0.662"],
        ["exp05_pca_boost", "HistGradientBoosting+PCA", "39 829", "79 027", "0.594"],
    ]
    add_table(slide, Inches(0.3), Inches(1.3), Inches(12.7), Inches(3.2), headers, rows,
              col_widths=[Inches(2.5), Inches(3.5), Inches(2.0), Inches(2.2), Inches(2.5)])

    items = [
        "Подбор гиперпараметров: RandomizedSearchCV (5-fold) только на train",
        "Выбор финальной модели: минимум RMSE на val → GradientBoostingRegressor",
        "На test: RMSE 64 019 (−8% vs RF из CP1), R² 0.733",
    ]
    add_bullet_list(slide, Inches(0.7), Inches(4.9), Inches(11.5), Inches(2.0),
                    items, font_size=16, color=LIGHT_GRAY)

    # ================================================================
    # SLIDE 6 — Интерпретируемость
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
                 "Интерпретируемость: Permutation Importance", font_size=32, bold=True, color=ACCENT)

    # Bar chart
    chart_data = CategoryChartData()
    chart_data.categories = [
        "conversions", "industry_vert.", "income_bracket",
        "creative_age", "clicks", "avg_session_dur",
        "campaign_obj.", "hour_of_day", "CTR", "creative_emot.",
    ]
    chart_data.add_series("RMSE increase", (76500, 7514, 4786, 661, 543, 296, 289, 277, 231, 226))

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.5), Inches(1.3),
        Inches(7.5), Inches(5.5), chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False

    plot = chart.plots[0]
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = ACCENT

    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(11)
    cat_axis.tick_labels.font.color.rgb = WHITE
    val_axis = chart.value_axis
    val_axis.tick_labels.font.size = Pt(10)
    val_axis.tick_labels.font.color.rgb = LIGHT_GRAY

    items = [
        "conversions — главный драйвер:\n  объём конверсий напрямую связан\n  с выручкой (но не утечка: это\n  наблюдаемый трафик, не cost)",
        "industry_vertical, income_bracket —\n  отрасль и доход аудитории\n  определяют маржинальность",
        "creative_age_days, clicks —\n  «свежесть» креатива и объём\n  трафика",
        "Метод: sklearn permutation_importance,\n  n_repeats=5, scoring=neg_RMSE, на val",
    ]
    add_bullet_list(slide, Inches(8.3), Inches(1.3), Inches(4.5), Inches(5.5),
                    items, font_size=14, color=LIGHT_GRAY)

    # ================================================================
    # SLIDE 7 — Covariate shift & ablation
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
                 "Сдвиг распределений и ablation", font_size=32, bold=True, color=ACCENT)

    items = [
        "Time-based split: val RMSE ≈ 27k, test RMSE ≈ 64k — разрыв из-за covariate shift",
        "Календарные признаки (quarter, date_month, date_dayofyear) — наибольший сдвиг",
        "Test-период попадает в «другой сезон» → модель экстраполирует",
        "",
        "Ablation: random split 70/15/15 vs time split →",
        "  Random split даёт заниженный RMSE (модель «подглядывает» будущее)",
        "  Time split — честная оценка обобщения на production-данных",
        "",
        "Вывод: time split — правильный выбор для задачи с временным порядком",
    ]
    add_bullet_list(slide, Inches(0.7), Inches(1.3), Inches(11.5), Inches(5.5),
                    items, font_size=17)

    # ================================================================
    # SLIDE 8 — Деплой
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
                 "Деплой: FastAPI + Streamlit + Docker", font_size=32, bold=True, color=ACCENT)

    # Architecture diagram as text
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.0),
                 "┌─────────────────────┐\n"
                 "│   Streamlit UI      │\n"
                 "│   :8501             │\n"
                 "│   форма ввода фичей │\n"
                 "└─────────┬───────────┘\n"
                 "          │ POST /predict\n"
                 "          ▼\n"
                 "┌─────────────────────┐\n"
                 "│   FastAPI  :8000    │\n"
                 "│   /health           │\n"
                 "│   /model            │\n"
                 "│   /predict          │\n"
                 "│   /predict_batch    │\n"
                 "└─────────┬───────────┘\n"
                 "          │ joblib.load\n"
                 "          ▼\n"
                 "┌─────────────────────┐\n"
                 "│  final_cp2.joblib   │\n"
                 "│  Pipeline: prep →   │\n"
                 "│  OHE → GBR         │\n"
                 "└─────────────────────┘",
                 font_size=12, color=LIGHT_GRAY, font_name="Courier New")

    items = [
        "docker compose up — поднимает оба сервиса",
        "Модель загружается 1 раз при старте (lifespan)",
        "Pydantic-схема: 33 признака, без leakage-колонок",
        "Swagger UI: localhost:8000/docs",
        "Streamlit UI: localhost:8501",
        "",
        "Препроцессинг не дублируется — всё внутри Pipeline",
        "OneHotEncoder(handle_unknown='ignore') → устойчив к новым категориям",
        "",
        "Тесты: 7 smoke-тестов через TestClient",
        "  /health → 200, /predict → 200 (число)",
        "  /predict (missing field) → 422",
        "  sanity: API == pipe.predict()",
    ]
    add_bullet_list(slide, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.5),
                    items, font_size=14, color=WHITE)

    # ================================================================
    # SLIDE 9 — Качество кода
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
                 "Качество кода и воспроизводимость", font_size=32, bold=True, color=ACCENT)

    headers = ["Критерий", "Реализация"]
    rows = [
        ["SEED", "42 — фиксирован в src/config.py, вызывается set_seed() в каждом ноутбуке"],
        ["Линтер", "ruff check src/ --line-length 120 — CI + локально"],
        ["Тесты", "pytest — 17 тестов (10 pipeline + 7 API), все зелёные"],
        ["Docker", "docker-compose.yml — два сервиса (api, ui), воспроизводимый образ"],
        ["requirements.txt", "Pin-версии всех зависимостей (sklearn 1.4.2, fastapi 0.115.12, ...)"],
        ["CI", ".github/workflows/ci.yml — ruff + pytest на каждый push"],
        ["Структура", "src/ (config, preprocessing, modeling, api/, ui/) + tests/ + notebooks/"],
        ["Документация", "README.md — быстрый старт, структура, деплой, примеры curl"],
    ]
    add_table(slide, Inches(0.3), Inches(1.3), Inches(12.7), Inches(5.0), headers, rows,
              col_widths=[Inches(2.8), Inches(9.9)])

    # ================================================================
    # SLIDE 10 — Итоги
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
                 "Итоги и что можно улучшить", font_size=32, bold=True, color=ACCENT)

    # Results summary
    add_text_box(slide, Inches(0.7), Inches(1.2), Inches(5.5), Inches(0.5),
                 "Результаты", font_size=22, bold=True, color=WHITE)

    items_left = [
        "Финальная модель: GradientBoostingRegressor",
        "Test RMSE: 64 019  |  R²: 0.733",
        "Улучшение vs CP1 RF: −8% RMSE",
        "Деплой: FastAPI + Streamlit + Docker",
        "17 тестов, ruff, CI — всё зелёное",
        "Отчёт: 8 разделов + PDF",
    ]
    add_bullet_list(slide, Inches(0.7), Inches(1.8), Inches(5.5), Inches(4.0),
                    items_left, font_size=17, color=WHITE)

    # What to improve
    add_text_box(slide, Inches(7.0), Inches(1.2), Inches(5.5), Inches(0.5),
                 "Что улучшить", font_size=22, bold=True, color=ORANGE)

    items_right = [
        "Квантильная регрессия — оценка неопределённости",
        "Монотонные ограничения в бустинге\n  (conversions ↑ → profit ↑)",
        "Переобучение на новых данных через API",
        "Мониторинг data drift в production",
        "Target encoding для high-cardinality фичей",
    ]
    add_bullet_list(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.0),
                    items_right, font_size=17, color=LIGHT_GRAY)

    # Bottom line
    add_text_box(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.8),
                 "Спасибо за внимание!  Вопросы?",
                 font_size=28, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)

    # ================================================================
    # Save
    # ================================================================
    out_path = "presentation/slides.pptx"
    prs.save(out_path)
    print(f"Saved → {out_path}")


if __name__ == "__main__":
    main()
